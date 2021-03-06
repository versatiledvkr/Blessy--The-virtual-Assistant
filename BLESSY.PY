import pyttsx3 
import speech_recognition as sr 
import datetime
import wikipedia 
import webbrowser
import os
import pyjokes
import pywhatkit as kit
import wolframalpha
import docx
import random
from fpdf import FPDF
from pptx import Presentation
import pyautogui
from playsound import playsound
from GoogleNews import GoogleNews


engine = pyttsx3.init('sapi5')
voices = engine.getProperty('voices')
# print(voices[1].id)
engine.setProperty('voice', voices[0].id)


def speak(audio):
    engine.say(audio)
    engine.runAndWait()


def wishMe():
    hour = int(datetime.datetime.now().hour)
    playsound("romantic.mp3")
    print("*****{ BLESSY-THE POWERFUL AI PERSONAL ASSISTANT }***** vamsi krishna's creation")
    speak("HI  i'm  BLESSY an AI personal assistant  made by vamsi krishna Please tell me how may i help you")       

def takeCommand():
    #It takes microphone input from the user and returns string output

    r = sr.Recognizer()
    with sr.Microphone() as source:
        print("listening")
        r.pause_threshold = 1
        audio = r.listen(source)

    try:
        print("recognising")    
        query = r.recognize_google(audio, language='eng-in')
        print(f"you said: {query}\n")

    except Exception as e:
        # print(e)    
        print("*")
        speak('say something else')  
        return "None"
    return query



def screenshot ():
    img = pyautogui.screenshot()
    img.save("B.png")

app = wolframalpha.Client("LPPQ5A-U9WPHJA6RY")


if __name__ == "__main__":
    wishMe()
    while True:
    # if 1:
        query = takeCommand().lower()

        if 'baby' in query:
            
            query = query.replace("wikipedia", "")
            results = wikipedia.summary(query, sentences=10)
            print(results)
            speak(results)

        if 'stop' in query or 'over' in query or 'bye' in query or 'quit' in query or 'see you' in query:
            f = "bye ", "ok bye ", "see you again ", "bye bye", "As your wish ", "tata bye "
            speak(random.choice(f))
            
            break

        elif "who am i" in query or "can you guess me" in query:
            jh = "if you are speaking then, definately you are a human", "You are vamsi", "You are a human", "I cant identify peoples with their vocies, may be you are vamsi or anybody with relation of vamsi"
            speak(random.choice(jh))


        elif 'hi' in query or "hello" in query:
            speak("hello there,how can i help with")

        elif 'how are you' in query or "how do you do" in query:
            speak("iam fine and i hope everything is alright from your side")

        elif "who created you" in query or "who made you" in query :
            speak("i was made by vamsi krishna")

        elif "I love you" in query:
            speak("I love you too")

        elif "when is your birthday" in query:
            speak(" are you palning for a suprise party oh just chill it's on 4th january")
        
        elif "what is your date of birth" in query:
            speak("4th january 2021")

        elif "what is your age" in query:
            speak ("i'm ageless i move up according to the generations")

        elif "your gender" in query:
            speak("i'm genderless but my voice seems to be like female")

        elif "who is god" in query or "who is your god" in query:
            speak("the super power which created you and helps you in all aspects in your life according to me you humans are the gods to A I world")

        elif 'what can you do' in query:
            speak(" i can do lot many things like setting alarms and reminders i can attend online classes and meetings i can send whatsapp messages and emails i can play songs and videos i can search anything from google i can answer all your question using my ai brain i can launch all your applications i can write word documents and create powerpoint presentaions i can control your pc such as taking screenshot and volume controls i can even shut down your pc i can even make calls i can calculate anything i can even say wether forecast ")
        
        elif "what are your hobbies" in query or "what is your goal" in query:
            speak("nothing just providing my service to you")

        elif "why you are so kind" in query:
            speak("because you gave me life")

        elif "thanks" in query or "thank you" in query:
            speak("it's my pleasure")

        elif 'open instagram' in query:
            webbrowser.open("instagram.com")
            speak("opening instagram")
            

        elif 'open vfma' in query:
            webbrowser.open("https://rajuvamsikrishna40.wixsite.com/mysite")
            speak("opening vfma")

        elif 'sing with me' in query or 'say along with me' in query:
            speak("as per your order")
            padu =takeCommand()
            speak("yes"+ padu)
            
        elif 'weather' in query:
            res = app.query(query)
            speak(next(res.results).text)
            print(next(res.results).text)

        elif 'calculate' in query:
            speak("what should i calculate")
            gh = takeCommand().lower()
            res = app.query(gh)
            speak(next(res.results).text)
            print(next(res.results).text)

        elif 'shutdown' in query:
            speak ("within how many seconds i should shut down the pc")
            kit.shutdown(int(input()))
            speak("shuting down")

        elif 'cancel shutdown' in query:
            speak("cancelling shutdown")
            kit.cancelShutdown()
        

        elif 'tell me a joke' in query:
            print(pyjokes.get_joke())
            speak(pyjokes.get_joke())
            
        elif 'the time' in query:
            strTime = datetime.datetime.now().strftime("%H:%M:%S")    
            speak(f"Sir, the time is {strTime}")

        elif 'send message' in query:
            speak ("what should i say")
            say= takeCommand()
            speak("please enter number and by what time i should send message")
            kit.sendwhatmsg("+91"+input(),f"{say}",int(input()),int(input()))
            speak('message sent')

        elif 'play song' in query:
            speak ("what song should i play")
            say= takeCommand() 
            kit.playonyt(f"{say}")
            speak("playing" + f"{say}")

        elif 'play video' in query:
            speak ("what video should i play")
            say= takeCommand() 
            kit.playonyt(f"{say}")
            speak("playing" + f"{say}")

        
        
        elif 'search' in query: 
            speak("what should i search")
            say= takeCommand() 
            kit.search(f"{say}") 
            speak("searching"+ f"{say}")
    

        elif 'create a word document' in query:
            speak("what should i write")
            say= takeCommand()
            doc = docx.Document()
            parag= doc.add_paragraph(f"{say}")
            parag.add_run(".This document was created by ")
            parag.add_run("blessy").bold=True
            doc.save("B.docx")
            os.system("start B.docx")

        elif 'create a pdf' in query:
            speak("what should i note")
            say=takeCommand()
            pdf = FPDF() 
            pdf.add_page()
            pdf.set_font("Arial",size=10)
            pdf.cell(200,100, txt= f"{say}")
            pdf.output("B.pdf")
            os.system("start B.pdf")

        elif 'create a presentation' in query:
            speak("what is the content")
            say= takeCommand()
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title= slide.shapes.title
            subtitle= slide.placeholders[1]
            title.text = f"{say}"
            subtitle.text = "Presentation by Blessy the A.I personal assistant made by vamsi krishna"
            prs.save("B.pptx")
            os.system("start B.pptx")

        elif 'start my meeting' in query:
            speak("starting your meeting")
            webbrowser.open("https://us04web.zoom.us/j/4411595128?pwd=ZlVRbXhZR21jQitJQThjcWNiN0lpUT09")

        elif 'join math a class'in query:
            speak("joining math a class")
            webbrowser.open("https://zoom.us/j/2654731902?pwd=eDVDSUFtWFJqVnY2YzdGRWYyMlBJQT09")

        elif 'join math b class' in query:
            speak("joining math b class")
            webbrowser.open("https://zoom.us/j/7619850243?pwd=LzFWV1dicGkyTFhMbVpwcEMvSnVDQT09")

        elif 'take screenshot' in query:
            screenshot()
            speak("screenshot captured")

        elif 'volume up' in query:
            pyautogui.press("volumeup")

        elif 'volume down' in query:
            pyautogui.press("volumedown")

        elif 'mute' in query:
            pyautogui.press("volumemute")

        elif 'turn on gaming mode' in query:
            pyautogui.hotkey("fn","f7")

        elif 'news' in query:
            googlenews = GoogleNews(lang='en', region='US')
        

        elif 'remember' in query:
            speak("what should i remember ")
            note = takeCommand()
            remember = open('data.txt','w')
            remember.write(note)
            remember.close()
            speak("i remember that"+ note)

        elif 'what are my reminders' in query:
            remember = open("data.txt","r").read()
            speak("you told me to remind that" + remember)
            
        elif 'send email' in query:
             speak("please enter sender's email id")
             sender= input()
             speak("please enter the password,it is safe don't worry")
             password= input()
             speak("what is the subject")
             subject= takeCommand()
             speak("what is the content")
             content= takeCommand()
             speak ("please enter the receiver email id")
             receiver= input()
             kit.send_mail(sender,password,subject,content,receiver)     

        
        elif 'alarm' in query:
            speak('Enter The Time :')
            time = input(": Enter The Time :")

            while True:
               Time_Ac = datetime.datetime.now()
               now = Time_Ac.strftime("%H:%M:%S")

               if now == time: 
                   speak("get up")
                   playsound("BLESSY tone.mp3")
                   speak ("alarm closed")

                
                
               elif now>time:
                    break
